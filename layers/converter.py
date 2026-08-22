# -*- coding: utf-8 -*-
"""LibreOffice转换封装；供上传、工具箱和Agent附件转换链路复用。"""

import os
import shutil
import subprocess
import threading
import time
import uuid
from enum import Enum
from typing import Optional

from pydantic import BaseModel

import config
from layers.file_processing.libreoffice import LibreOfficeProcessor
from layers.file_processing.models import FileProcessingRequest, FileTaskType
from layers.file_processing.runtime import (
    get_file_processor_registry,
    register_processor_once,
)
from layers.pdf_text import extract_pdf_page_text
from utils.logger import get_logger


logger = get_logger("converter")
_conversion_lock = threading.Lock()
_pdf_conversion_lock = threading.Lock()


class ConversionStatus(str, Enum):
    SUCCESS = "SUCCESS"
    FAILED = "FAILED"
    TIMEOUT = "TIMEOUT"


class ConversionResult(BaseModel):
    success: bool
    status: ConversionStatus
    output_path: Optional[str] = None
    converted_from_format: str = ""
    converted_to_format: str = ""
    error_type: str = ""
    error_msg: str = ""


def _convert_file_impl(source_path: str, target_format: str) -> ConversionResult:
    """Convert one local file through headless soffice under a process-wide lock."""
    started_at = time.perf_counter()
    output_dir = ""
    source_ext = os.path.splitext(source_path or "")[1].lower()
    target = (target_format or "").lower().lstrip(".")
    try:
        if not source_path or not os.path.isfile(source_path):
            return _failed("待转换文件不存在", source_ext, target, "invalid_source")
        if target not in {"docx", "pdf"}:
            return _failed("不支持的转换目标格式", source_ext, target, "invalid_target")
        max_bytes = max(0, config.MAX_CONVERSION_FILE_SIZE_MB) * 1024 * 1024
        if os.path.getsize(source_path) > max_bytes:
            return _failed("文件超过转换大小限制", source_ext, target, "file_too_large")

        soffice_path = _resolve_soffice_path()
        if not soffice_path:
            return _failed(
                "服务器未安装或未配置LibreOffice，暂时无法转换该格式",
                source_ext,
                target,
                "not_configured",
            )

        output_dir = os.path.join(
            os.path.dirname(source_path),
            "conversion_%s" % uuid.uuid4().hex,
        )
        os.makedirs(output_dir, exist_ok=False)
        command = [
            soffice_path,
            "--headless",
            "--convert-to",
            target,
            "--outdir",
            output_dir,
            source_path,
        ]
        with _conversion_lock:
            completed = subprocess.run(
                command,
                capture_output=True,
                check=False,
                timeout=max(1, config.CONVERSION_TIMEOUT_SECONDS),
            )
        if completed.returncode != 0:
            _cleanup_directory(output_dir)
            return _failed("LibreOffice转换失败", source_ext, target, "process_failed")

        expected_path = os.path.join(
            output_dir,
            "%s.%s" % (os.path.splitext(os.path.basename(source_path))[0], target),
        )
        if not os.path.isfile(expected_path):
            _cleanup_directory(output_dir)
            return _failed("LibreOffice未生成转换文件", source_ext, target, "output_missing")

        _log_conversion(source_ext, target, "success", started_at)
        return ConversionResult(
            success=True,
            status=ConversionStatus.SUCCESS,
            output_path=expected_path,
            converted_from_format=source_ext.lstrip("."),
            converted_to_format=target,
        )
    except subprocess.TimeoutExpired:
        _cleanup_directory(output_dir)
        _log_conversion(source_ext, target, "timeout", started_at)
        return ConversionResult(
            success=False,
            status=ConversionStatus.TIMEOUT,
            converted_from_format=source_ext.lstrip("."),
            converted_to_format=target,
            error_type="timeout",
            error_msg="文档转换超时，请稍后重试",
        )
    except Exception as exc:
        _cleanup_directory(output_dir)
        logger.warning(
            "文档转换异常：source_ext=%s target=%s error_type=%s",
            source_ext,
            target,
            type(exc).__name__,
        )
        return _failed("文档转换失败", source_ext, target, type(exc).__name__)


def convert_pdf_to_office(source_path: str, target_format: str) -> ConversionResult:
    """Best-effort PDF content reconstruction without OCR."""
    started_at = time.perf_counter()
    target = (target_format or "").lower().lstrip(".")
    output_dir = ""
    try:
        if not source_path or not os.path.isfile(source_path):
            return _failed("待转换文件不存在", ".pdf", target, "invalid_source")
        if target not in {"docx", "xlsx", "pptx"}:
            return _failed("不支持的PDF转换目标", ".pdf", target, "invalid_target")
        max_bytes = max(0, config.MAX_CONVERSION_FILE_SIZE_MB) * 1024 * 1024
        if os.path.getsize(source_path) > max_bytes:
            return _failed("文件超过转换大小限制", ".pdf", target, "file_too_large")
        output_dir = os.path.join(
            os.path.dirname(source_path),
            "conversion_%s" % uuid.uuid4().hex,
        )
        os.makedirs(output_dir, exist_ok=False)
        output_path = os.path.join(
            output_dir,
            "%s.%s" % (os.path.splitext(os.path.basename(source_path))[0], target),
        )
        with _pdf_conversion_lock:
            if target == "docx":
                _pdf_to_docx(source_path, output_path)
            elif target == "xlsx":
                _pdf_to_xlsx(source_path, output_path)
            else:
                _pdf_to_pptx(source_path, output_path)
        if not os.path.isfile(output_path):
            _cleanup_directory(output_dir)
            return _failed("未生成转换文件", ".pdf", target, "output_missing")
        _log_conversion(".pdf", target, "success", started_at)
        return ConversionResult(
            success=True,
            status=ConversionStatus.SUCCESS,
            output_path=output_path,
            converted_from_format="pdf",
            converted_to_format=target,
        )
    except Exception as exc:
        _cleanup_directory(output_dir)
        error_type = "encrypted_pdf" if "password" in str(exc).lower() else type(exc).__name__
        logger.warning(
            "PDF反向转换异常：target=%s error_type=%s",
            target,
            error_type,
        )
        return _failed("PDF内容提取或重建失败", ".pdf", target, error_type)


def _pdf_to_docx(source_path: str, output_path: str) -> None:
    import pdfplumber
    from docx import Document

    document = Document()
    with pdfplumber.open(source_path) as pdf:
        for page_index, page in enumerate(pdf.pages):
            if page_index:
                document.add_page_break()
            text = extract_pdf_page_text(page)
            for line in text.splitlines():
                document.add_paragraph(line)
    document.save(output_path)


def _pdf_to_xlsx(source_path: str, output_path: str) -> None:
    import pdfplumber
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.remove(workbook.active)
    with pdfplumber.open(source_path) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            sheet = workbook.create_sheet("Page %s" % page_index)
            row_index = 1
            tables = page.extract_tables() or []
            if tables:
                for table in tables:
                    for row in table:
                        for column_index, value in enumerate(row or [], start=1):
                            sheet.cell(row=row_index, column=column_index, value=value or "")
                        row_index += 1
                    row_index += 1
            else:
                for line in extract_pdf_page_text(page).splitlines():
                    sheet.cell(row=row_index, column=1, value=line)
                    row_index += 1
    if not workbook.sheetnames:
        workbook.create_sheet("Page 1")
    workbook.save(output_path)


def _pdf_to_pptx(source_path: str, output_path: str) -> None:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    pdf = fitz.open(source_path)
    try:
        for page_index, page in enumerate(pdf):
            image_path = os.path.join(
                os.path.dirname(output_path),
                "page_%s.png" % (page_index + 1),
            )
            page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).save(image_path)
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                image_path,
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
            os.remove(image_path)
    finally:
        pdf.close()
    presentation.save(output_path)


def cleanup_conversion_output(output_path: str) -> None:
    """Remove a successful conversion artifact and its private output directory."""
    if not output_path:
        return
    output_dir = os.path.dirname(output_path)
    if os.path.basename(output_dir).startswith("conversion_"):
        _cleanup_directory(output_dir)
        return
    try:
        if os.path.isfile(output_path):
            os.remove(output_path)
    except OSError as exc:
        logger.warning("转换产物清理失败：error_type=%s", type(exc).__name__)


def _resolve_soffice_path() -> str:
    configured = (config.LIBREOFFICE_PATH or "").strip()
    if configured and os.path.isfile(configured):
        return configured
    discovered = shutil.which("soffice")
    return discovered or ""


def _cleanup_directory(path: str) -> None:
    if not path:
        return
    try:
        shutil.rmtree(path, ignore_errors=False)
    except FileNotFoundError:
        return
    except OSError as exc:
        logger.warning("转换临时目录清理失败：error_type=%s", type(exc).__name__)


def _failed(
    message: str,
    source_ext: str = "",
    target: str = "",
    error_type: str = "failed",
) -> ConversionResult:
    return ConversionResult(
        success=False,
        status=ConversionStatus.FAILED,
        converted_from_format=source_ext.lstrip("."),
        converted_to_format=target,
        error_type=error_type,
        error_msg=message,
    )


def _log_conversion(source_ext: str, target: str, status: str, started_at: float) -> None:
    logger.info(
        "文档转换完成：source_ext=%s target=%s status=%s elapsed_ms=%s",
        source_ext,
        target,
        status,
        int((time.perf_counter() - started_at) * 1000),
    )


_libreoffice_processor = LibreOfficeProcessor(
    conversion_delegate=_convert_file_impl,
    cleanup_delegate=cleanup_conversion_output,
    max_size_bytes=max(0, config.MAX_CONVERSION_FILE_SIZE_MB) * 1024 * 1024,
)
register_processor_once(_libreoffice_processor)


def convert_file(source_path: str, target_format: str) -> ConversionResult:
    """通过统一注册表裁决后调用既有LibreOffice稳定链路。"""
    source_ext = os.path.splitext(source_path or "")[1].lower()
    target = (target_format or "").lower().lstrip(".")
    request = FileProcessingRequest(
        task_type=FileTaskType.CONVERT,
        source_paths=[source_path] if source_path else [],
        source_format=source_ext,
        target_format=target,
        max_input_size_bytes=max(0, config.MAX_CONVERSION_FILE_SIZE_MB) * 1024 * 1024,
        max_output_size_bytes=max(0, config.MAX_CONVERSION_FILE_SIZE_MB) * 1024 * 1024,
    )
    try:
        processor, _ = get_file_processor_registry().resolve(request)
    except LookupError:
        return _convert_file_impl(source_path, target_format)
    result = processor.execute(request)
    if not result.success:
        return ConversionResult(
            success=False,
            status=ConversionStatus(result.status.value),
            converted_from_format=source_ext.lstrip("."),
            converted_to_format=target,
            error_type=result.error_type,
            error_msg=result.error_message,
        )
    quality = processor.validate_output(request, result)
    if not quality.passed or quality.artifact is None:
        processor.cleanup(request, result)
        return _failed(
            "转换产物质量检查未通过",
            source_ext,
            target,
            quality.issues[0].code if quality.issues else "quality_check_failed",
        )
    return ConversionResult(
        success=True,
        status=ConversionStatus.SUCCESS,
        output_path=quality.artifact.output_path,
        converted_from_format=source_ext.lstrip("."),
        converted_to_format=target,
    )
