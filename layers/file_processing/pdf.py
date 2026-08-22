# -*- coding: utf-8 -*-
"""统一PDF处理器；保留pdfplumber、pypdf、fitz各自已验证的职责。"""

import os
import shutil
import threading
import time
import uuid
from typing import List, Optional

import fitz
import pdfplumber
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter

import config
from layers.file_processing.base import FileProcessor
from layers.file_processing.models import (
    FileArtifact,
    FileProcessingRequest,
    FileProcessingResult,
    FileProcessingStatus,
    FileTaskType,
    ProcessorCapability,
    QualityCheckResult,
    QualityIssue,
    QualityProfile,
)
from layers.file_processing.quality import FileQualityChecker
from layers.file_processing.runtime import register_processor_once
from layers.pdf_text import extract_pdf_page_text
from utils.logger import get_logger


logger = get_logger("converter")
_pdf_processing_lock = threading.Lock()
_MIME_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "png": "image/png",
}
_QUALITY_PROFILES = {
    "pdf": QualityProfile.PDF,
    "docx": QualityProfile.DOCX,
    "xlsx": QualityProfile.XLSX,
    "pptx": QualityProfile.PPTX,
    "png": QualityProfile.PNG,
}


class PdfProcessor(FileProcessor):
    name = "pdf"
    adapter_version = "1"

    def __init__(self, max_size_bytes: int) -> None:
        self._max_size_bytes = max_size_bytes
        self._quality_checker = FileQualityChecker()

    def capabilities(self) -> List[ProcessorCapability]:
        common = {
            "processor_name": self.name,
            "source_formats": ["pdf"],
            "asynchronous": True,
            "max_size_bytes": self._max_size_bytes,
            "requires_external_binary": False,
        }
        return [
            ProcessorCapability(
                capability_id="pdf.extract_text",
                target_formats=[],
                task_types=[FileTaskType.EXTRACT_TEXT],
                output_mime_types=["text/plain"],
                knowledge_base_eligible=True,
                quality_profile=QualityProfile.PDF,
                **common,
            ),
            ProcessorCapability(
                capability_id="pdf.extract_tables",
                target_formats=[],
                task_types=[FileTaskType.EXTRACT_TABLES],
                output_mime_types=["application/json"],
                knowledge_base_eligible=True,
                quality_profile=QualityProfile.PDF,
                **common,
            ),
            ProcessorCapability(
                capability_id="pdf.render_pages",
                target_formats=["png"],
                task_types=[FileTaskType.RENDER_PAGES],
                output_mime_types=[_MIME_TYPES["png"]],
                knowledge_base_eligible=False,
                quality_profile=QualityProfile.PNG,
                **common,
            ),
            ProcessorCapability(
                capability_id="pdf.merge",
                target_formats=["pdf"],
                task_types=[FileTaskType.MERGE],
                output_mime_types=[_MIME_TYPES["pdf"]],
                knowledge_base_eligible=False,
                quality_profile=QualityProfile.PDF,
                **common,
            ),
            ProcessorCapability(
                capability_id="pdf.split",
                target_formats=["pdf"],
                task_types=[FileTaskType.SPLIT],
                output_mime_types=[_MIME_TYPES["pdf"]],
                knowledge_base_eligible=False,
                quality_profile=QualityProfile.PDF,
                **common,
            ),
            *[
                ProcessorCapability(
                    capability_id="pdf.to_%s" % target,
                    target_formats=[target],
                    task_types=[FileTaskType.CONVERT],
                    output_mime_types=[_MIME_TYPES[target]],
                    knowledge_base_eligible=True,
                    quality_profile=_QUALITY_PROFILES[target],
                    **common,
                )
                for target in ("docx", "xlsx", "pptx")
            ],
        ]

    def supports(self, request: FileProcessingRequest) -> bool:
        if request.source_format != "pdf":
            return False
        if request.task_type == FileTaskType.CONVERT:
            return request.target_format in {"docx", "xlsx", "pptx"}
        return request.task_type in {
            FileTaskType.EXTRACT_TEXT,
            FileTaskType.EXTRACT_TABLES,
            FileTaskType.RENDER_PAGES,
            FileTaskType.MERGE,
            FileTaskType.SPLIT,
        }

    def execute(self, request: FileProcessingRequest) -> FileProcessingResult:
        try:
            validation = self._validate_sources(request)
            if validation is not None:
                return validation
            if request.task_type == FileTaskType.EXTRACT_TEXT:
                return self._extract_text(request.source_paths[0])
            if request.task_type == FileTaskType.EXTRACT_TABLES:
                return self._extract_tables(request.source_paths[0])
            if request.task_type == FileTaskType.RENDER_PAGES:
                return self._render_pages(request.source_paths[0], request.output_dir or "")
            if request.task_type == FileTaskType.MERGE:
                return self._merge(request.source_paths, request.output_path or "")
            if request.task_type == FileTaskType.SPLIT:
                return self._split(
                    request.source_paths[0],
                    request.output_dir or "",
                    request.max_pages,
                )
            if request.task_type == FileTaskType.CONVERT:
                return self._convert(request.source_paths[0], request.target_format)
            return self._failed("unsupported_task", "不支持的PDF任务")
        except ValueError as exc:
            self._cleanup_failed_request(request)
            error_type = str(exc) if str(exc) == "encrypted_pdf" else "invalid_pdf"
            return self._failed(error_type, "PDF已加密" if error_type == "encrypted_pdf" else "PDF文件损坏")
        except Exception as exc:
            self._cleanup_failed_request(request)
            if request.task_type == FileTaskType.CONVERT:
                error_type = (
                    "encrypted_pdf"
                    if "password" in str(exc).lower()
                    else type(exc).__name__
                )
                logger.warning(
                    "PDF反向转换异常：target=%s error_type=%s",
                    request.target_format,
                    error_type,
                )
                return self._failed(error_type, "PDF内容提取或重建失败")
            logger.warning(
                "PDF处理异常：task_type=%s target=%s error_type=%s",
                request.task_type.value,
                request.target_format,
                type(exc).__name__,
            )
            return self._failed("invalid_pdf", "PDF文件损坏或无法解析")

    def validate_output(
        self,
        request: FileProcessingRequest,
        result: FileProcessingResult,
    ) -> QualityCheckResult:
        if not result.success:
            return QualityCheckResult(passed=False)
        if request.task_type in {FileTaskType.EXTRACT_TEXT, FileTaskType.EXTRACT_TABLES}:
            return QualityCheckResult(passed=True)
        issues: List[QualityIssue] = []
        for artifact in result.artifacts:
            checked = self._quality_checker.validate(
                artifact,
                _QUALITY_PROFILES[artifact.file_format],
                max_size_bytes=request.max_output_size_bytes,
                minimum_pages=1 if artifact.file_format in {"pdf", "pptx", "png"} else 0,
                minimum_paragraphs=1 if artifact.file_format == "docx" else 0,
                minimum_worksheets=1 if artifact.file_format == "xlsx" else 0,
            )
            issues.extend(checked.issues)
        return QualityCheckResult(
            passed=not issues and bool(result.artifacts),
            artifact=result.artifacts[0] if not issues and result.artifacts else None,
            issues=issues,
        )

    def cleanup(
        self,
        request: FileProcessingRequest,
        result: FileProcessingResult,
    ) -> None:
        paths = [artifact.output_path for artifact in result.artifacts]
        conversion_dirs = {
            os.path.dirname(path)
            for path in paths
            if os.path.basename(os.path.dirname(path)).startswith("conversion_")
        }
        for path in paths:
            if os.path.isfile(path) and os.path.dirname(path) not in conversion_dirs:
                os.remove(path)
        for directory in conversion_dirs:
            shutil.rmtree(directory)
        if request.output_dir and os.path.isdir(request.output_dir):
            try:
                os.rmdir(request.output_dir)
            except OSError as exc:
                logger.warning(
                    "PDF临时目录清理异常：error_type=%s",
                    type(exc).__name__,
                )

    @staticmethod
    def _cleanup_failed_request(request: FileProcessingRequest) -> None:
        try:
            if request.output_path and os.path.isfile(request.output_path):
                os.remove(request.output_path)
            if request.output_dir and os.path.isdir(request.output_dir):
                shutil.rmtree(request.output_dir)
        except OSError as exc:
            logger.warning("PDF失败产物清理异常：error_type=%s", type(exc).__name__)

    def _validate_sources(
        self,
        request: FileProcessingRequest,
    ) -> Optional[FileProcessingResult]:
        if not request.source_paths or any(
            not path or not os.path.isfile(path) for path in request.source_paths
        ):
            return self._failed("invalid_source", "待处理PDF不存在")
        if request.max_input_size_bytes > 0 and any(
            os.path.getsize(path) > request.max_input_size_bytes
            for path in request.source_paths
        ):
            return self._failed("file_too_large", "文件超过转换大小限制")
        return None

    def _extract_text(self, source_path: str) -> FileProcessingResult:
        texts = []
        with pdfplumber.open(source_path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                text = extract_pdf_page_text(page)
                if text.strip():
                    texts.append(text)
        return self._success(text="\n\n".join(texts), page_count=page_count)

    def _extract_tables(self, source_path: str) -> FileProcessingResult:
        tables: List[List[List[Optional[str]]]] = []
        with pdfplumber.open(source_path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                for table in page.extract_tables() or []:
                    tables.append([
                        [None if value is None else str(value) for value in row or []]
                        for row in table
                    ])
        return self._success(tables=tables, page_count=page_count)

    def _render_pages(self, source_path: str, output_dir: str) -> FileProcessingResult:
        os.makedirs(output_dir, exist_ok=True)
        artifacts = []
        document = fitz.open(source_path)
        try:
            for page_index, page in enumerate(document, start=1):
                output_path = os.path.join(output_dir, "page_%s.png" % page_index)
                page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).save(output_path)
                artifacts.append(self._artifact(output_path, "png"))
        finally:
            document.close()
        return self._success(artifacts=artifacts, page_count=len(artifacts))

    def _merge(self, source_paths: List[str], output_path: str) -> FileProcessingResult:
        writer = PdfWriter()
        try:
            for source_path in source_paths:
                reader = self._reader(source_path)
                for page in reader.pages:
                    writer.add_page(page)
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            with open(output_path, "wb") as output_file:
                writer.write(output_file)
            return self._success(
                artifacts=[self._artifact(output_path, "pdf")],
                page_count=len(writer.pages),
            )
        finally:
            writer.close()

    def _split(
        self,
        source_path: str,
        output_dir: str,
        max_pages: int,
    ) -> FileProcessingResult:
        reader = self._reader(source_path)
        page_count = len(reader.pages)
        if max_pages > 0 and page_count > max_pages:
            return self._failed("too_many_pages", "PDF页数超过拆分上限", page_count)
        os.makedirs(output_dir, exist_ok=True)
        artifacts = []
        for index, page in enumerate(reader.pages, start=1):
            output_path = os.path.join(output_dir, "page_%s.pdf" % index)
            writer = PdfWriter()
            try:
                writer.add_page(page)
                with open(output_path, "wb") as output_file:
                    writer.write(output_file)
            finally:
                writer.close()
            artifacts.append(self._artifact(output_path, "pdf"))
        return self._success(artifacts=artifacts, page_count=page_count)

    def _convert(self, source_path: str, target_format: str) -> FileProcessingResult:
        started_at = time.perf_counter()
        output_dir = os.path.join(
            os.path.dirname(source_path),
            "conversion_%s" % uuid.uuid4().hex,
        )
        os.makedirs(output_dir, exist_ok=False)
        output_path = os.path.join(
            output_dir,
            "%s.%s" % (os.path.splitext(os.path.basename(source_path))[0], target_format),
        )
        try:
            with _pdf_processing_lock:
                if target_format == "docx":
                    self._to_docx(source_path, output_path)
                elif target_format == "xlsx":
                    self._to_xlsx(source_path, output_path)
                elif target_format == "pptx":
                    self._to_pptx(source_path, output_path)
                else:
                    return self._failed("invalid_target", "不支持的PDF转换目标")
            if not os.path.isfile(output_path):
                shutil.rmtree(output_dir)
                return self._failed("output_missing", "未生成转换文件")
            logger.info(
                "文档转换完成：source_ext=.pdf target=%s status=success elapsed_ms=%s",
                target_format,
                int((time.perf_counter() - started_at) * 1000),
            )
            return self._success(artifacts=[self._artifact(output_path, target_format)])
        except Exception:
            shutil.rmtree(output_dir, ignore_errors=True)
            raise

    @staticmethod
    def _to_docx(source_path: str, output_path: str) -> None:
        document = Document()
        with pdfplumber.open(source_path) as pdf:
            for page_index, page in enumerate(pdf.pages):
                if page_index:
                    document.add_page_break()
                for line in extract_pdf_page_text(page).splitlines():
                    document.add_paragraph(line)
        document.save(output_path)

    @staticmethod
    def _to_xlsx(source_path: str, output_path: str) -> None:
        workbook = Workbook()
        workbook.remove(workbook.active)
        with pdfplumber.open(source_path) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                sheet = workbook.create_sheet("Page %s" % page_index)
                row_index = 1
                tables = page.extract_tables() or []
                if tables:
                    for table in tables:
                        for row in table:
                            for column_index, value in enumerate(row or [], start=1):
                                sheet.cell(row=row_index, column=column_index, value=value or "")
                            row_index += 1
                        row_index += 1
                else:
                    for line in extract_pdf_page_text(page).splitlines():
                        sheet.cell(row=row_index, column=1, value=line)
                        row_index += 1
        if not workbook.sheetnames:
            workbook.create_sheet("Page 1")
        workbook.save(output_path)

    @staticmethod
    def _to_pptx(source_path: str, output_path: str) -> None:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        blank_layout = presentation.slide_layouts[6]
        document = fitz.open(source_path)
        try:
            for page_index, page in enumerate(document, start=1):
                image_path = os.path.join(
                    os.path.dirname(output_path),
                    "page_%s.png" % page_index,
                )
                page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).save(image_path)
                slide = presentation.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    image_path,
                    0,
                    0,
                    width=presentation.slide_width,
                    height=presentation.slide_height,
                )
                os.remove(image_path)
        finally:
            document.close()
        presentation.save(output_path)

    @staticmethod
    def _reader(path: str) -> PdfReader:
        reader = PdfReader(path)
        if reader.is_encrypted:
            raise ValueError("encrypted_pdf")
        return reader

    def _artifact(self, output_path: str, file_format: str) -> FileArtifact:
        return FileArtifact(
            output_path=output_path,
            file_format=file_format,
            mime_type=_MIME_TYPES[file_format],
            engine_name=self.name,
            engine_version=self.adapter_version,
        )

    @staticmethod
    def _success(
        artifacts: Optional[List[FileArtifact]] = None,
        text: str = "",
        tables: Optional[List[List[List[Optional[str]]]]] = None,
        page_count: int = 0,
    ) -> FileProcessingResult:
        return FileProcessingResult(
            success=True,
            status=FileProcessingStatus.SUCCESS,
            artifacts=artifacts or [],
            text=text,
            tables=tables or [],
            page_count=page_count,
        )

    @staticmethod
    def _failed(
        error_type: str,
        error_message: str,
        page_count: int = 0,
    ) -> FileProcessingResult:
        return FileProcessingResult(
            success=False,
            status=FileProcessingStatus.FAILED,
            page_count=page_count,
            error_type=error_type,
            error_message=error_message,
        )


def get_pdf_processing_lock() -> threading.Lock:
    return _pdf_processing_lock


pdf_processor = PdfProcessor(
    max_size_bytes=max(0, config.MAX_CONVERSION_FILE_SIZE_MB) * 1024 * 1024
)
register_processor_once(pdf_processor)
