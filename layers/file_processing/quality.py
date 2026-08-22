# -*- coding: utf-8 -*-
"""文件处理产物的通用质量检查。"""

import mimetypes
import os
import zipfile
from typing import List, Tuple

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from layers.file_processing.models import (
    FileArtifact,
    QualityCheckResult,
    QualityIssue,
    QualityProfile,
)


_PROFILE_MIME_TYPES = {
    QualityProfile.TEXT: {"text/plain", "text/markdown"},
    QualityProfile.PDF: {"application/pdf"},
    QualityProfile.DOCX: {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    },
    QualityProfile.XLSX: {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    },
    QualityProfile.PPTX: {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    },
}


class FileQualityChecker:
    """验证产物可交付性；任一问题都会阻止结果进入文件库。"""

    def validate(
        self,
        artifact: FileArtifact,
        profile: QualityProfile,
        max_size_bytes: int = 0,
        minimum_pages: int = 0,
        minimum_paragraphs: int = 0,
        minimum_worksheets: int = 0,
    ) -> QualityCheckResult:
        issues: List[QualityIssue] = []
        path = artifact.output_path
        if not os.path.isfile(path):
            return self._failed("output_missing", "处理产物不存在")
        size_bytes = os.path.getsize(path)
        artifact.size_bytes = size_bytes
        if size_bytes <= 0:
            issues.append(QualityIssue(code="output_empty", message="处理产物为空"))
        if max_size_bytes > 0 and size_bytes > max_size_bytes:
            issues.append(QualityIssue(code="output_too_large", message="处理产物超过大小限制"))
        issues.extend(self._validate_mime(path, artifact, profile))
        if not issues:
            issues.extend(
                self._reopen_and_measure(
                    artifact,
                    profile,
                    minimum_pages,
                    minimum_paragraphs,
                    minimum_worksheets,
                )
            )
        return QualityCheckResult(
            passed=not issues,
            artifact=artifact if not issues else None,
            issues=issues,
        )

    def validate_cleanup(self, paths: List[str]) -> QualityCheckResult:
        leftovers = [path for path in paths if path and os.path.exists(path)]
        if not leftovers:
            return QualityCheckResult(passed=True)
        return self._failed("temporary_files_remain", "处理器仍有临时文件未清理")

    def _validate_mime(
        self,
        path: str,
        artifact: FileArtifact,
        profile: QualityProfile,
    ) -> List[QualityIssue]:
        expected = _PROFILE_MIME_TYPES[profile]
        if artifact.mime_type not in expected:
            return [QualityIssue(code="mime_mismatch", message="声明的MIME与目标格式不符")]
        if profile == QualityProfile.PDF:
            with open(path, "rb") as source:
                if source.read(5) != b"%PDF-":
                    return [QualityIssue(code="format_mismatch", message="文件内容不是有效PDF")]
        elif profile in {QualityProfile.DOCX, QualityProfile.XLSX, QualityProfile.PPTX}:
            if not zipfile.is_zipfile(path):
                return [QualityIssue(code="format_mismatch", message="Office文件容器无效")]
        else:
            guessed = mimetypes.guess_type(path)[0]
            if guessed and guessed not in expected:
                return [QualityIssue(code="format_mismatch", message="文本扩展名与目标格式不符")]
        return []

    def _reopen_and_measure(
        self,
        artifact: FileArtifact,
        profile: QualityProfile,
        minimum_pages: int,
        minimum_paragraphs: int,
        minimum_worksheets: int,
    ) -> List[QualityIssue]:
        try:
            text_samples: List[str] = []
            if profile == QualityProfile.TEXT:
                with open(artifact.output_path, "r", encoding="utf-8") as source:
                    text_samples.append(source.read())
            elif profile == QualityProfile.PDF:
                reader = PdfReader(artifact.output_path)
                artifact.page_count = len(reader.pages)
                text_samples.extend((page.extract_text() or "") for page in reader.pages)
            elif profile == QualityProfile.DOCX:
                document = Document(artifact.output_path)
                artifact.paragraph_count = len(document.paragraphs)
                text_samples.extend(item.text for item in document.paragraphs)
            elif profile == QualityProfile.XLSX:
                workbook = load_workbook(artifact.output_path, read_only=True, data_only=True)
                try:
                    artifact.worksheet_count = len(workbook.worksheets)
                    for sheet in workbook.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            text_samples.extend(str(value) for value in row if value is not None)
                finally:
                    workbook.close()
            else:
                presentation = Presentation(artifact.output_path)
                artifact.page_count = len(presentation.slides)
                for slide in presentation.slides:
                    text_samples.extend(
                        shape.text for shape in slide.shapes if hasattr(shape, "text")
                    )
            issues = self._minimum_issues(
                artifact,
                minimum_pages,
                minimum_paragraphs,
                minimum_worksheets,
            )
            if self._contains_invalid_unicode(text_samples):
                issues.append(QualityIssue(code="text_corrupted", message="文本包含无效替换字符"))
            return issues
        except (OSError, ValueError, UnicodeError, zipfile.BadZipFile):
            return [QualityIssue(code="reopen_failed", message="处理产物无法重新打开")]
        except Exception:
            return [QualityIssue(code="reopen_failed", message="处理产物结构校验失败")]

    @staticmethod
    def _minimum_issues(
        artifact: FileArtifact,
        minimum_pages: int,
        minimum_paragraphs: int,
        minimum_worksheets: int,
    ) -> List[QualityIssue]:
        values: List[Tuple[str, int, int]] = [
            ("page_count_too_small", artifact.page_count, minimum_pages),
            ("paragraph_count_too_small", artifact.paragraph_count, minimum_paragraphs),
            ("worksheet_count_too_small", artifact.worksheet_count, minimum_worksheets),
        ]
        return [
            QualityIssue(code=code, message="处理产物结构少于最低预期")
            for code, actual, minimum in values
            if minimum > 0 and actual < minimum
        ]

    @staticmethod
    def _contains_invalid_unicode(values: List[str]) -> bool:
        return any("\ufffd" in value for value in values)

    @staticmethod
    def _failed(code: str, message: str) -> QualityCheckResult:
        return QualityCheckResult(
            passed=False,
            issues=[QualityIssue(code=code, message=message)],
        )
